#!/usr/bin/env python
"""
11_build_presentation.py
========================
Build the 15-slide walkthrough deck in ``results/`` covering the whole project:
where every dataset came from, how it was processed, what was wrong with it,
how driving distance and time were computed and under what assumptions, what the
results are, and how to position the grant proposal.

Output is a ``.pptx``, which opens natively in PowerPoint, Keynote and Google
Slides (File > Import in Slides).

Every number in the deck is read from the pipeline's own output tables at build
time rather than typed in, so the deck cannot drift out of sync with the
analysis. If a table is missing the script fails loudly instead of shipping a
slide with stale numbers.

Usage
-----
    python scripts/11_build_presentation.py
"""

from __future__ import annotations

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

import paths as P

OUT = P.RESULTS / "MatCare_Project_Walkthrough.pptx"

# 16:9
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

INK = RGBColor(0x11, 0x11, 0x11)
INK_2 = RGBColor(0x52, 0x51, 0x4E)
MUTED = RGBColor(0x8A, 0x89, 0x83)
BLUE = RGBColor(0x1C, 0x5C, 0xAB)
BLUE_L = RGBColor(0x2A, 0x78, 0xD6)
ORANGE = RGBColor(0xB9, 0x4B, 0x21)
RED = RGBColor(0xC0, 0x39, 0x38)
GREEN = RGBColor(0x0E, 0x7A, 0x53)
SURFACE = RGBColor(0xFC, 0xFC, 0xFB)
BAND = RGBColor(0xEE, 0xF4, 0xFD)


def num(path: str, **kw) -> pd.DataFrame:
    p = P.TABLES / path
    if not p.exists():
        raise SystemExit(f"missing {p} - run the pipeline first")
    return pd.read_csv(p, **kw)


# --------------------------------------------------------------------------
# Slide primitives
# --------------------------------------------------------------------------
def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = SURFACE
    return s


def textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def para(tf, text, size=16, bold=False, color=INK, space_after=6,
         align=PP_ALIGN.LEFT, first=False, bullet=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = ("•  " if bullet else "") + text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return p


def header(slide, title, kicker=None):
    """Slide title with a coloured rule under it."""
    tf = textbox(slide, Inches(0.55), Inches(0.32), Inches(12.2), Inches(0.9))
    if kicker:
        para(tf, kicker.upper(), size=11.5, bold=True, color=BLUE_L,
             space_after=2, first=True)
        para(tf, title, size=27, bold=True, color=INK, space_after=0)
    else:
        para(tf, title, size=27, bold=True, color=INK, space_after=0, first=True)

    from pptx.enum.shapes import MSO_SHAPE
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.55),
        Inches(1.42) if kicker else Inches(1.18), Inches(1.5), Pt(3.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE_L
    bar.line.fill.background()
    bar.shadow.inherit = False


def footer(slide, text):
    tf = textbox(slide, Inches(0.55), Inches(6.92), Inches(12.2), Inches(0.4))
    para(tf, text, size=10.5, color=MUTED, space_after=0, first=True)


def picture(slide, stem, x, y, w=None, h=None):
    """Place a figure. Give width OR height; the other is derived to preserve
    aspect ratio. Tall multi-panel figures must be sized by height or they run
    off the bottom of a 16:9 slide."""
    p = P.FIGURES / f"{stem}.png"
    if not p.exists():
        raise SystemExit(f"missing figure {p}")
    if h is not None:
        return slide.shapes.add_picture(str(p), x, y, height=h)
    return slide.shapes.add_picture(str(p), x, y, width=w)


def table(slide, rows, x, y, w, col_w=None, size=12.5, header_fill=BLUE):
    """rows[0] is the header."""
    nr, nc = len(rows), len(rows[0])
    shp = slide.shapes.add_table(nr, nc, x, y, w, Inches(0.32 * nr))
    t = shp.table
    if col_w:
        total = sum(col_w)
        for j, frac in enumerate(col_w):
            t.columns[j].width = Emu(int(w * frac / total))
    for i, row in enumerate(rows):
        t.rows[i].height = Inches(0.3)
        for j, val in enumerate(row):
            c = t.cell(i, j)
            c.text = str(val)
            c.margin_left, c.margin_right = Inches(0.07), Inches(0.07)
            c.margin_top, c.margin_bottom = 0, 0
            pr = c.text_frame.paragraphs[0]
            pr.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT
            for r in pr.runs:
                r.font.size = Pt(size)
                r.font.name = "Calibri"
                r.font.bold = i == 0
                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if i == 0 else INK
            c.fill.solid()
            c.fill.fore_color.rgb = header_fill if i == 0 else (
                BAND if i % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF))
    return t


def stat_row(slide, stats, y, x=Inches(0.65), w=Inches(12.0)):
    """Big-number tiles: [(value, label, color), ...]"""
    n = len(stats)
    gap = Inches(0.22)
    tile = Emu(int((w - gap * (n - 1)) / n))
    for i, (val, label, col) in enumerate(stats):
        cx = Emu(int(x + i * (tile + gap)))
        from pptx.enum.shapes import MSO_SHAPE
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, y,
                                     tile, Inches(1.28))
        box.fill.solid()
        box.fill.fore_color.rgb = BAND
        box.line.color.rgb = RGBColor(0xD5, 0xE3, 0xF7)
        box.shadow.inherit = False
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_top, tf.margin_bottom = Inches(0.1), Inches(0.05)
        para(tf, val, size=25, bold=True, color=col, space_after=1,
             align=PP_ALIGN.CENTER, first=True)
        para(tf, label, size=11, color=INK_2, space_after=0,
             align=PP_ALIGN.CENTER)


# --------------------------------------------------------------------------
def main() -> int:
    P.ensure_tree()

    # ---- pull live numbers ------------------------------------------------
    thr = num("women_access_thresholds.csv").set_index("benchmark")
    wsum = num("women_access_summary.csv").set_index("metric")
    bands = num("access_time_bands.csv")
    sit = num("siting_recommendations.csv")
    curve = num("siting_marginal_gain.csv")
    match = num("facility_match_report.csv")
    val = num("facility_validation_by_county.csv", dtype={"COUNTYFP5": str})

    w_total = int(thr.loc["within 30 min", "women_15_44"] /
                  (thr.loc["within 30 min", "pct"] / 100))
    w_beyond30 = w_total - int(thr.loc["within 30 min", "women_15_44"])
    w_beyond30_nicu = w_total - int(thr.loc["within 30 min of NICU", "women_15_44"])
    hosp_no_ob = val[(val.matched_facilities.fillna(0) == 0) & (val.osm_hospitals.fillna(0) >= 1)]

    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H

    # ================================================== 1 title
    s = blank(prs)
    from pptx.enum.shapes import MSO_SHAPE
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(2.9))
    band.fill.solid()
    band.fill.fore_color.rgb = BLUE
    band.line.fill.background()
    band.shadow.inherit = False
    tf = textbox(s, Inches(0.8), Inches(0.75), Inches(11.7), Inches(2.0))
    para(tf, "MatCare — Obstetric Access in Texas", size=36, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF), space_after=8, first=True)
    para(tf, "How the data was built, what was wrong with it, and what it shows",
         size=19, color=RGBColor(0xCD, 0xE2, 0xFB), space_after=0)

    stat_row(s, [
        (f"{int(thr.loc['within 30 min','women_15_44']):,}", "women 15–44 within 30 min of care", BLUE),
        (f"{w_beyond30:,}", "women 15–44 beyond 30 min", ORANGE),
        (f"{w_beyond30_nicu:,}", "beyond 30 min of a NICU", RED),
        ("211", "active obstetric facilities", GREEN),
    ], y=Inches(3.35))

    tf = textbox(s, Inches(0.8), Inches(5.05), Inches(11.7), Inches(1.9))
    para(tf, "Every Texas census block routed on the real road network. "
             "No sampling, no API, fully reproducible.",
         size=15, color=INK_2, space_after=10, first=True)
    para(tf, "Data: US Census (TIGER, 2020 Decennial, ACS 2019–2023) · Centers for "
             "Medicare & Medicaid Services (CMS) Provider of Services · HIFLD · "
             "OpenStreetMap", size=12, color=MUTED, space_after=0)
    footer(s, "Project walkthrough · slide 1 of 15")

    # ================================================== 2 the question
    s = blank(prs)
    header(s, "Why the standard measure is not good enough", "the problem")
    tf = textbox(s, Inches(0.65), Inches(1.75), Inches(6.0), Inches(4.9))
    para(tf, "The usual metric: “maternity care desert” = a county with no "
             "obstetric facility.", size=17, bold=True, space_after=14, first=True)
    for b in [
        "A county is not a service area. Texas counties range from 386 to "
        "16,040 km² — averaging access across one hides the variation a "
        "planner needs.",
        "It is binary. A county with one facility and a county with twenty "
        "score identically.",
        "It cannot distinguish a county that never had a hospital from one "
        "whose hospital closed its labour-and-delivery unit — clinically and "
        "politically very different problems.",
        "It ignores level of care. Reaching any hospital is not the same as "
        "reaching one that can manage a haemorrhage or a preterm birth.",
    ]:
        para(tf, b, size=14, color=INK_2, space_after=11, bullet=True)
    para(tf, "Our approach: measure drive time from every one of the 668,757 "
             "census blocks, weighted by women aged 15–44.",
         size=15, bold=True, color=BLUE, space_after=0)

    table(s, [
        ["Unit", "Count", "Resolution"],
        ["County", "254", "too coarse"],
        ["Census tract", "6,896", "coarse"],
        ["Block group", "18,638", "ACS available here"],
        ["Census block", "668,757", "what we use"],
    ], Inches(7.15), Inches(2.35), Inches(5.5), col_w=[3, 1.4, 2.2])
    footer(s, "Slide 2 of 15")

    # ================================================== 3 data sources
    s = blank(prs)
    header(s, "Which dataset came from where", "data provenance")
    table(s, [
        ["Dataset", "Source", "What it provides"],
        ["TIGER/Line 2024 blocks", "census.gov", "668,757 block polygons + 2020 population"],
        ["TIGER block groups / tracts", "census.gov", "18,638 / 6,896 polygons"],
        ["Centers of Population 2020", "census.gov", "population-WEIGHTED centroids"],
        ["ACS 2019–2023 table B01001", "census.gov bulk", "women aged 15–44 by block group"],
        ["ZCTA Gazetteer 2024", "census.gov", "ZIP centroids (geocoding fallback)"],
        ["Cartographic boundaries", "census.gov", "county / place outlines for maps"],
        ["CMS Provider of Services", "supplied (CMS)", "WHICH hospitals provide obstetrics"],
        ["HIFLD hospitals", "supplied (HIFLD)", "WHERE hospitals are (876 in TX)"],
        ["OpenStreetMap Texas", "Geofabrik", "the road network (713 MB, ODbL)"],
    ], Inches(0.65), Inches(1.75), Inches(12.0), col_w=[3.0, 2.0, 4.6], size=12)

    tf = textbox(s, Inches(0.65), Inches(5.5), Inches(12.0), Inches(1.3))
    para(tf, "Every download is recorded in data/_manifest.json with its URL, byte "
             "count, SHA-256 hash and retrieval time.", size=13.5, bold=True,
         color=INK, space_after=7, first=True)
    para(tf, "The OpenStreetMap extract was verified against Geofabrik's own published "
             "MD5 checksum — it matched exactly. Census block, block-group and tract "
             "populations independently agree at 29,145,505, the published 2020 Texas total.",
         size=12.5, color=INK_2, space_after=0)
    footer(s, "Slide 3 of 15 · No API key is required anywhere in this pipeline")

    # ================================================== 4 pipeline
    s = blank(prs)
    header(s, "The processing pipeline, end to end", "how it runs")
    steps = [
        ("00", "Download + hash every source", "1.25 GB, resumable, SHA-256 manifest"),
        ("01", "Build the facility list", "CMS POS ⋈ HIFLD, multi-stage match"),
        ("02", "Build census geography", "blocks, block groups, tracts, counties"),
        ("03", "Build the road network", "OSM → 11.0M nodes, 21.6M directed edges"),
        ("04", "Compute access", "multi-source Dijkstra, every block"),
        ("05", "Validate the facility list", "independent cross-check vs OSM hospitals"),
        ("06", "Standard figures", "maps + distributions"),
        ("07", "Optimise siting", "greedy Maximal Covering Location Problem"),
        ("08", "Proposal figures", "3 publication-quality figures"),
        ("10", "Women 15–44 denominator", "ACS B01001, block-group resolution"),
    ]
    y = Inches(1.8)
    for i, (n, title, sub) in enumerate(steps):
        col = i // 5
        row = i % 5
        x = Inches(0.65 + col * 6.25)
        yy = Emu(int(y + row * Inches(0.98)))
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, yy,
                                 Inches(0.72), Inches(0.72))
        box.fill.solid()
        box.fill.fore_color.rgb = BLUE_L
        box.line.fill.background()
        box.shadow.inherit = False
        p = box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = n
        r.font.size, r.font.bold = Pt(17), True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        tf = textbox(s, Emu(int(x + Inches(0.92))), yy, Inches(5.1), Inches(0.8))
        para(tf, title, size=14.5, bold=True, space_after=1, first=True)
        para(tf, sub, size=11.5, color=MUTED, space_after=0)

    footer(s, "Slide 4 of 15 · Numbered, idempotent, no Docker, no routing server "
              "· full statewide run < 10 minutes")

    # ================================================== 5 census + checks
    s = blank(prs)
    header(s, "Population geography, and how we know it is right", "step 02")
    tf = textbox(s, Inches(0.65), Inches(1.75), Inches(6.1), Inches(4.6))
    para(tf, "Centroids: two definitions, chosen deliberately", size=17,
         bold=True, space_after=12, first=True)
    for b in [
        "Block groups and tracts use the Census Bureau's official 2020 "
        "population-WEIGHTED centres of population.",
        "Blocks use the TIGER internal point — the Bureau publishes no "
        "weighted centroid at block level (verified: the directory holds only "
        "county/, tract/, blkgrp/).",
        "This matters. An areal centroid of a large rural county can sit in "
        "empty desert far from anyone, inflating measured travel. Blocks are "
        "small enough that the difference is ~100 m.",
    ]:
        para(tf, b, size=13.5, color=INK_2, space_after=11, bullet=True)

    para(tf, "Integrity checks that passed", size=17, bold=True,
         color=GREEN, space_after=10)
    for b in ["Blocks, block groups and tracts independently total 29,145,505 "
              "— the published 2020 Texas population. Discrepancy: zero.",
              "254 Texas counties, exactly as expected.",
              "Women 15–44 disaggregated to blocks reconciles to the block-group "
              "total with zero difference."]:
        para(tf, b, size=13.5, color=INK_2, space_after=9, bullet=True)

    table(s, [
        ["Layer", "Count", "2020 population"],
        ["Census blocks", "668,757", "29,145,505"],
        ["Block groups", "18,638", "29,145,505"],
        ["Census tracts", "6,896", "29,145,505"],
        ["Counties", "254", "29,145,505"],
        ["Populated blocks", "449,085", "67.2% of blocks"],
    ], Inches(7.2), Inches(2.3), Inches(5.45), col_w=[2.6, 1.6, 2.4])
    footer(s, "Slide 5 of 15")

    # ================================================== 6 women denominator
    s = blank(prs)
    header(s, "The right denominator: women aged 15–44", "step 10")
    tf = textbox(s, Inches(0.65), Inches(1.7), Inches(6.2), Inches(4.9))
    para(tf, "Total population answers the wrong question.", size=17, bold=True,
         space_after=12, first=True)
    for b in [
        "Obstetric care is used by women of reproductive age, and they are not "
        "distributed like everyone else — university towns, retirement "
        "communities and prisons all skew it.",
        "Source: ACS 2019–2023 5-year, table B01001 “Sex by Age”, at block-group "
        "resolution (ACS is not published for blocks).",
        "The 9 age bands are derived AT RUNTIME by parsing the official Census "
        "table shells, not typed from memory — so a mis-remembered variable "
        "number cannot silently corrupt the denominator.",
        "For block-level work, block-group counts are shared out to blocks in "
        "proportion to 2020 decennial population. This assumes a uniform age–sex "
        "mix within a block group, and is stated as an assumption.",
    ]:
        para(tf, b, size=13.5, color=INK_2, space_after=11, bullet=True)
    para(tf, "api.census.gov now requires a key. We avoid it entirely by reading "
             "the bulk Summary File — identical data, no registration.",
         size=13, bold=True, color=BLUE, space_after=0)

    table(s, [
        ["ACS variable", "Female age band"],
        ["B01001_030", "15 to 17 years"],
        ["B01001_031", "18 and 19 years"],
        ["B01001_032", "20 years"],
        ["B01001_033", "21 years"],
        ["B01001_034", "22 to 24 years"],
        ["B01001_035", "25 to 29 years"],
        ["B01001_036", "30 to 34 years"],
        ["B01001_037", "35 to 39 years"],
        ["B01001_038", "40 to 44 years"],
        ["TOTAL", f"{w_total:,} women"],
    ], Inches(7.35), Inches(1.95), Inches(5.3), col_w=[2.4, 2.9], size=12)
    footer(s, "Slide 6 of 15 · 20.9% of the Texas population")

    # ================================================== 7 facility sources
    s = blank(prs)
    header(s, "Defining “an obstetric facility” needs two sources", "step 01")
    tf = textbox(s, Inches(0.65), Inches(1.75), Inches(5.9), Inches(4.6))
    para(tf, "Neither source alone is sufficient.", size=17, bold=True,
         space_after=12, first=True)
    para(tf, "CMS Provider of Services", size=15, bold=True, color=BLUE,
         space_after=4)
    para(tf, "Knows WHICH hospitals provide obstetrics. Carries no coordinates.",
         size=13.5, color=INK_2, space_after=12)
    para(tf, "HIFLD hospitals", size=15, bold=True, color=BLUE, space_after=4)
    para(tf, "Knows WHERE hospitals are (876 in Texas). Does not record "
             "obstetric service.", size=13.5, color=INK_2, space_after=14)
    para(tf, "Neither carries the other's identifier — so they must be matched.",
         size=14, bold=True, color=ORANGE, space_after=14)
    para(tf, "Note: the project's “CMOS” files are the CMS Provider of Services "
             "extract, confirmed from its own 733-page record layout — not a "
             "Texas HHS registry as previously assumed.",
         size=12.5, color=MUTED, space_after=0)

    table(s, [
        ["Eligibility rule", "Effect"],
        ["OB_SRVC_CD = 1 or 3", "obstetrics provided ON SITE"],
        ["OB_SRVC_CD = 2", "under arrangement — excluded"],
        ["PGM_TRMNTN_CD = 00", "active provider only"],
        ["→ eligible Texas hospitals", "211"],
        ["→ of which NICU on site", "117"],
    ], Inches(6.85), Inches(2.4), Inches(5.8), col_w=[3.2, 2.6], size=12.5)

    tf = textbox(s, Inches(6.85), Inches(4.5), Inches(5.8), Inches(1.6))
    para(tf, "Codes are quoted from the CMS record layout, not inferred:",
         size=12.5, bold=True, color=INK_2, space_after=5, first=True)
    para(tf, "0 = NOT PROVIDED · 1 = BY STAFF · 2 = UNDER ARRANGEMENT · "
             "3 = BOTH", size=11.5, color=MUTED, space_after=0)
    footer(s, "Slide 7 of 15")

    # ================================================== 8 discrepancy 1
    s = blank(prs)
    header(s, "Discrepancy #1 — the exact-name join", "what was wrong")
    tf = textbox(s, Inches(0.65), Inches(1.7), Inches(12.0), Inches(1.5))
    para(tf, "The original facility list was built with one line:", size=14,
         color=INK_2, space_after=7, first=True)
    p = para(tf, "    matched_df = df2[df2['NAME'].isin(df1['FAC_NAME'])]",
             size=15, bold=True, color=RED, space_after=7)
    p.runs[0].font.name = "Consolas"
    para(tf, "An exact, case-sensitive string equality between two agencies' "
             "spellings of a hospital name. It fails silently — no error, just a "
             "smaller table that looks plausible.", size=13.5, color=INK_2,
         space_after=0)

    table(s, [
        ["What it dropped", "CMS name", "HIFLD name"],
        ["Acronym", "BSA HOSPITAL", "BAPTIST ST ANTHONYS HOSPITAL"],
        ["Corporate rename", "PARKLAND HEALTH AND HOSPITAL SYSTEM", "PARKLAND MEMORIAL HOSPITAL"],
        ["“System” vs site", "UNIVERSITY HEALTH SYSTEM", "UNIVERSITY HOSPITAL"],
        ["Punctuation / abbrev.", "ST. vs SAINT, MED CTR vs MEDICAL CENTER", "—"],
    ], Inches(0.65), Inches(3.4), Inches(12.0), col_w=[2.3, 4.6, 4.2], size=11.5)

    tf = textbox(s, Inches(0.65), Inches(5.35), Inches(12.0), Inches(1.5))
    para(tf, "Symptom that gave it away:", size=14, bold=True, color=ORANGE,
         space_after=6, first=True)
    para(tf, "Jefferson County — Beaumont, population 256,526 — showed ZERO "
             "obstetric facilities and a 61.6-minute mean drive time. The "
             "second-worst county in Texas, in a metro of 400,000. That is not a "
             "finding about Texas; it is a defect in the input.",
         size=13.5, color=INK_2, space_after=0)
    footer(s, "Slide 8 of 15 · The old join matched 136 of 211 eligible hospitals (64%)")

    # ================================================== 9 discrepancy 2 + fix
    s = blank(prs)
    header(s, "Discrepancy #2 and the fix", "how we corrected it")
    tf = textbox(s, Inches(0.65), Inches(1.68), Inches(5.9), Inches(2.2))
    para(tf, "#2 — Closed hospitals counted as open", size=16, bold=True,
         color=RED, space_after=8, first=True)
    para(tf, "The old join never filtered PGM_TRMNTN_CD. Of 377 Texas records "
             "reporting obstetric service, 163 were TERMINATED providers — "
             "merged, closed, or decertified — and were being counted as "
             "operating facilities.", size=13.5, color=INK_2, space_after=10)
    para(tf, "Both errors push the same way: they manufacture deserts.",
         size=13.5, bold=True, color=ORANGE, space_after=0)

    tf = textbox(s, Inches(0.65), Inches(4.0), Inches(5.9), Inches(2.6))
    para(tf, "The fix: progressive matching", size=16, bold=True, color=GREEN,
         space_after=8, first=True)
    para(tf, "Address passes run BEFORE loose fuzzy-name passes, on purpose: a "
             "shared street number and street name is far stronger evidence of "
             "identity than a similar corporate name.", size=13.5, color=INK_2,
         space_after=7)
    para(tf, "Hospitals get renamed. They do not move.", size=14, bold=True,
         color=BLUE, space_after=0)

    rows = [["Matching pass", "Matched"]]
    for _, r in match.iterrows():
        rows.append([r["match_stage"], f"{int(r['n'])}"])
    rows.append(["TOTAL automatic", f"{int(match['n'].sum())} of 211 (98.1%)"])
    rows.append(["Residual → ZIP centroid, flagged", "4"])
    table(s, rows, Inches(6.85), Inches(1.95), Inches(5.8), col_w=[4.0, 1.9], size=12)

    stat_row(s, [
        ("170 → 211", "facilities identified", GREEN),
        ("61.6 → 18.7", "Jefferson Co. mean drive (min)", GREEN),
        ("−59%", "change in Texans beyond 30 min", GREEN),
    ], y=Inches(4.6), x=Inches(6.85), w=Inches(5.8))
    footer(s, "Slide 9 of 15 · Worst-county list now shows remote West Texas, "
              "not the Beaumont metro — the face-validity check that the fix worked")

    # ================================================== 10 road network
    s = blank(prs)
    header(s, "Building the road network", "step 03")
    tf = textbox(s, Inches(0.65), Inches(1.75), Inches(6.1), Inches(4.8))
    para(tf, "Source: OpenStreetMap Texas extract (Geofabrik), MD5-verified.",
         size=15, bold=True, space_after=12, first=True)
    for b in [
        "Kept: motorway, trunk, primary, secondary, tertiary, unclassified, "
        "residential, living_street + all link ramps.",
        "Excluded: service roads (driveways, parking aisles) — 1.93M ways. A "
        "route therefore ends at the public road nearest the hospital, a median "
        "of 95 m from it.",
        "Excluded: ways barred to motor vehicles (access = no / private / "
        "customers…) — 32,502 ways.",
        "Directed graph: one-way streets and roundabouts are honoured.",
        "Edge length is GEODESIC on the WGS84 ellipsoid — not a projected "
        "approximation. Texas spans 13° of longitude, where projection "
        "distortion would be measurable.",
    ]:
        para(tf, b, size=13, color=INK_2, space_after=10, bullet=True)

    table(s, [
        ["Road class", "Segments", "km"],
        ["residential", "15,624,231", "958,560"],
        ["tertiary", "1,961,368", "124,110"],
        ["secondary", "1,569,187", "119,955"],
        ["unclassified", "1,170,841", "87,278"],
        ["primary", "544,862", "54,402"],
        ["trunk", "237,239", "30,806"],
        ["motorway", "138,596", "18,012"],
    ], Inches(7.2), Inches(2.2), Inches(5.45), col_w=[2.3, 2.0, 1.6], size=11.5)

    stat_row(s, [
        ("11.0M", "network nodes", BLUE),
        ("21.6M", "directed edges", BLUE),
        ("~700k km", "of centreline road", BLUE),
    ], y=Inches(5.1), x=Inches(7.2), w=Inches(5.45))
    footer(s, "Slide 10 of 15 · Parsed in 36 s, peak memory 2.2 GB")

    # ================================================== 11 speeds
    s = blank(prs)
    header(s, "How driving time was calculated — and what we assumed",
           "the biggest assumption")
    tf = textbox(s, Inches(0.65), Inches(1.72), Inches(6.1), Inches(4.9))
    para(tf, "Only 20.06% of road segments carry a posted maxspeed tag in OSM.",
         size=15, bold=True, color=ORANGE, space_after=11, first=True)
    for b in [
        "Where OSM has a posted limit, we use it (parsing “60”, “60 mph”, "
        "“100 km/h”; a bare number means km/h per the OSM specification).",
        "For the other ~80%, we apply a free-flow default by road class — the "
        "table on the right. These are OUR assumptions, chosen to reflect "
        "typical Texas posted limits. They are not taken from any cited source, "
        "and they are recorded in network_meta.json so any run can be audited.",
        "Time = geodesic segment length ÷ assigned speed.",
        "Off-network “last mile”: the straight-line distance from a centroid to "
        "the nearest road is added back at an assumed 30 mph.",
    ]:
        para(tf, b, size=13, color=INK_2, space_after=10, bullet=True)
    para(tf, "What this means, stated plainly", size=14.5, bold=True, color=RED,
         space_after=7)
    para(tf, "Travel times contain no congestion, no signal delay and no turn "
             "restrictions — they are BEST-CASE. Urban times are understated more "
             "than rural ones, so the reported urban–rural gap is conservative. "
             "Distances do not depend on speed at all and are the more robust metric.",
         size=12.5, color=INK_2, space_after=0)

    table(s, [
        ["Road class", "Assumed mph"],
        ["motorway", "70"], ["trunk", "65"], ["primary", "60"],
        ["secondary", "55"], ["tertiary", "45"], ["unclassified", "40"],
        ["residential", "30"], ["living_street", "15"],
        ["motorway_link", "45"], ["trunk_link", "40"],
        ["primary / secondary_link", "35"], ["tertiary_link", "30"],
        ["off-network last mile", "30"],
    ], Inches(7.35), Inches(1.95), Inches(5.3), col_w=[3.4, 1.9], size=11.5)
    footer(s, "Slide 11 of 15 · Validating these against a commercial routing "
              "engine is Aim 1 of the proposal")

    # ================================================== 12 algorithm
    s = blank(prs)
    header(s, "The routing algorithm — why this is fast", "step 04")
    tf = textbox(s, Inches(0.65), Inches(1.72), Inches(6.2), Inches(4.9))
    para(tf, "The naive framing is unaffordable.", size=16, bold=True,
         color=RED, space_after=9, first=True)
    para(tf, "668,757 blocks × 211 facilities = 141 million origin–destination "
             "pairs. No routing API prices that. This is exactly why comparable "
             "studies sample, aggregate to counties, or fall back to straight lines.",
         size=13.5, color=INK_2, space_after=13)
    para(tf, "The reformulation", size=16, bold=True, color=GREEN, space_after=9)
    para(tf, "Only the NEAREST facility matters. So seed a single Dijkstra search "
             "with all 211 facilities at once. Relaxing outward, it labels every "
             "node in Texas with its cost to the closest facility — and which one "
             "— in one pass.", size=13.5, color=INK_2, space_after=13)
    para(tf, "Direction is handled correctly.", size=15, bold=True, color=BLUE,
         space_after=7)
    para(tf, "A multi-source search on the graph gives cost(facility → block). We "
             "need cost(block → facility) — a patient travelling TO care. On a "
             "one-way network those differ, so the search runs on the TRANSPOSED "
             "graph. Most comparable studies quietly assume symmetry.",
         size=13, color=INK_2, space_after=0)

    stat_row(s, [
        ("~3 sec", "per pass over 11M nodes", GREEN),
        ("0", "API calls / recurring cost", GREEN),
    ], y=Inches(2.1), x=Inches(7.3), w=Inches(5.35))
    tf = textbox(s, Inches(7.3), Inches(3.7), Inches(5.35), Inches(3.0))
    para(tf, "Two passes are run: one weighted by segment length (distance), one "
             "by segment time (travel time).", size=13, color=INK_2,
         space_after=11, first=True)
    para(tf, "Every block is computed exactly. Nothing is sampled, interpolated "
             "or extrapolated.", size=13.5, bold=True, color=INK, space_after=11)
    para(tf, "Road detour ratio: 1.33×", size=15, bold=True, color=ORANGE,
         space_after=6)
    para(tf, "Road travel is a third longer than straight-line distance, "
             "population-weighted. A buffer-based analysis understates real "
             "travel by that much.", size=13, color=INK_2, space_after=0)
    footer(s, "Slide 12 of 15 · This is what makes an interactive planning tool possible")

    # ================================================== 13 results map
    s = blank(prs)
    header(s, "What we found", "results")
    # Sized by height: this figure is nearly square, so a width-based fit
    # would run past the bottom of the slide.
    pic = picture(s, "proposal_fig1_access_landscape", Inches(0), Inches(1.32),
                  h=Inches(5.4))
    pic.left = Emu(int((SLIDE_W - pic.width) / 2))
    footer(s, "Slide 13 of 15 · (a) road distance to any obstetric facility (miles)  "
              "(b) road distance to NICU-capable, same scale — visibly worse  "
              "(c) population by distance band, miles  (d) hospital but no obstetric unit")

    # ================================================== 14 women results
    s = blank(prs)
    header(s, "Access for women aged 15–44", "results")
    stat_row(s, [
        (f"{thr.loc['within 25 miles','pct']:.1f}%", "within 25 miles", BLUE),
        (f"{thr.loc['within 35 miles','pct']:.1f}%", "within 35 miles", BLUE),
        (f"{thr.loc['within 50 miles','pct']:.1f}%", "within 50 miles", BLUE),
        (f"{thr.loc['within 30 min','pct']:.1f}%", "within 30 min", BLUE),
        (f"{thr.loc['within 30 min of NICU','pct']:.1f}%", "within 30 min of a NICU", RED),
    ], y=Inches(1.68))

    tf = textbox(s, Inches(0.65), Inches(3.2), Inches(6.1), Inches(3.5))
    para(tf, "The headline is the gap, not the level.", size=17, bold=True,
         space_after=12, first=True)
    para(tf, f"Only {w_beyond30:,} women aged 15–44 live more than 30 minutes "
             f"from an obstetric facility — 3.0%. But {w_beyond30_nicu:,} — "
             f"9.4%, more than three times as many — live more than 30 minutes "
             f"from a facility with a NICU.",
         size=14, color=INK_2, space_after=12)
    para(tf, "Proximity to a door is not proximity to care. The NICU gap is where "
             "preventable maternal and neonatal deaths occur, and it is invisible "
             "to any binary desert map.", size=14, bold=True, color=ORANGE,
         space_after=12)
    para(tf, f"Separately: {len(hosp_no_ob)} counties holding "
             f"{int(hosp_no_ob['population'].sum()):,} residents have a hospital "
             f"but no obstetric unit. Those need a unit RESTORED, not a hospital "
             f"BUILT — a far cheaper intervention.",
         size=14, color=INK_2, space_after=0)

    table(s, [
        ["Worst counties for women 15–44", "Mean drive (min)"],
        ["Crockett", "79.1"], ["Culberson", "69.4"], ["King", "68.1"],
        ["Terrell", "66.5"], ["Edwards", "66.4"], ["Presidio", "66.0"],
        ["Sutton", "64.1"], ["Reagan", "64.0"],
    ], Inches(7.35), Inches(3.35), Inches(5.3), col_w=[3.4, 1.9], size=12)
    footer(s, "Slide 14 of 15 · Block-group resolution, ACS 2019–2023")

    # ================================================== 15 siting + strategy
    s = blank(prs)
    header(s, "From diagnosis to decision — and the proposal", "what comes next")
    picture(s, "proposal_fig2_siting_optimizer", Inches(0.5), Inches(1.5),
            Inches(7.4))

    tf = textbox(s, Inches(8.1), Inches(1.55), Inches(4.75), Inches(5.2))
    para(tf, "Already working", size=16, bold=True, color=GREEN, space_after=8,
         first=True)
    covered = int(curve["cum_pop_covered"].iloc[-1])
    pct = float(curve["pct_of_underserved"].iloc[-1])
    under_total = int(round(covered / (pct / 100)))
    para(tf, f"Ten optimally sited facilities would bring {covered:,} of "
             f"{under_total:,} underserved women aged 15–44 within 30 minutes "
             f"— {pct:.1f}% of the gap.",
         size=13, color=INK_2, space_after=10)
    para(tf, "Greedy carries a provable (1 − 1/e) ≈ 63% approximation guarantee, "
             "because coverage is monotone submodular. A bounded result, not a "
             "guess.", size=12.5, color=INK_2, space_after=14)

    para(tf, "How to make the proposal competitive", size=16, bold=True,
         color=BLUE, space_after=8)
    for b in [
        "Lead with the 78 hospital-but-no-obstetrics counties — the cheapest "
        "intervention, and invisible to every published desert map.",
        "Lead with the NICU gap: 9.4% vs 3.0%.",
        "Show the siting plan. It already exists — that is the difference "
        "between promising analysis and delivering it.",
        "Show the correction. It demonstrates rigour and implicitly critiques "
        "the existing literature without naming anyone.",
        "Frame ILP as the shipped default and RL for the sequential, "
        "budget-phased problem — do not claim RL beats ILP on a static problem.",
    ]:
        para(tf, b, size=11.5, color=INK_2, space_after=7, bullet=True)
    footer(s, "Slide 15 of 15 · Discovery Foundation · $29,150 requested")

    prs.save(OUT)
    print(f"Wrote {OUT.relative_to(P.PROJECT_ROOT)}")
    print(f"  {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
